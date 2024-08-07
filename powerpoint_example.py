#1 import library

from pptx import Presentation

#2 create presentataion 

presentation = Presentation()

#3 create slide

#3.1 layout selection
slide_layout=presentation.slide_layouts[1] #title and content layout


#3.2 slide creation
slide=presentation.slides.add_slide(slide_layout)


#4 add title
slide.placeholders[0].text = "My Title"



#5 add content

slide.placeholders[1].text = "Hello World\nBye"

#6 save presentation
presentation.save("test.pptx")