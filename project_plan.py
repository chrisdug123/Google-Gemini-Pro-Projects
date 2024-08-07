import google.generativeai as palm
from Markdown2docx import Markdown2docx
from openpyxl import Workbook
from pptx import Presentation

def generate( api_key,project_charter):
    palm.configure(api_key= api_key)
    defaults = {
        'temperature' : 0.7,
        'candidate_count':1,
        'top_k':40,
        'top_p': 0.95,
        'max_output_tokens': 1024,
        'stop_sequences':[],
        #'safety_settings':[{"category":"HARM_CATEGORY_HARASSMENT","threshold":1},{"category":"HARM_CATEGORY_HATE_SPEECH","threshold":1},{"category":"HARM_CATEGORY_SEXUALLY_EXPLICIT","threshold":1},{"category":"HARM_CATEGORY_DANEROUS_CONTENT","threshold":2}]
    }
    # prompting!

    #product_idea=input()
    prompt = f"""Act as a senior project manager with experience in software project. Project Plan based on the Project Charter below:
    
    {project_charter}

    --

    END OF PROJECT CHARTER
    
    The output of the Project Plan needs to be in tabular format. The columns are:

    1) Task Name: Description of the task
    2) Duration: Time required to complete task.
    3) Dependencies: Other tasks that must be completed before this task.
    4) Status: Current status of the task (e.g., Not Started, In Progressm Completed).
    5) Resources: Tools, team memers, software, infrastructure, etc., required for the task.

    Example of the output:
    Task Name|Duration|Dependencies|Status|Resources
    Creating Plan|1 day|None|Not Started|Leadership team
    Documentation|3 days|Creating Plan|Not Started|Development team, internal software

    Don't add any additional content or notes after you finish listing the tasks. Add at least 10 tasks and no extra content after. Now please write the Project Plan in tabular format:

    """
   
    response = palm.generate_text(
        **defaults,
        prompt=prompt )

    # save that as a Word document
    #my_file = open('risk_management.md', 'w')
    #my_file.write(response.result)
    #my_file.close()

    #word_file=Markdown2docx('risk_management')
    #word_file.eat_soup()
    #word_file.write_html()
    #word_file.save()
    rows = response.result.split('\n')

    spreadsheet =[]

    for row in rows:
        split_row = row.split('|')
        spreadsheet.append(split_row)

    excel_save(spreadsheet)
    powerpoint_save(spreadsheet)
    


    return response.result

def excel_save(spreadsheet):
    wb = Workbook()
    ws = wb.active

    for row in spreadsheet:
        ws.append(row)

    wb.save('Project_Plan.xlsx')


def powerpoint_save(spreadsheet):
    presentation = Presentation()
    slide_layout=presentation.slide_layouts[1] #title and content layout

    spreadsheet=spreadsheet[5:]

    for row in spreadsheet:
        if len(row)<5:
            continue
        slide = presentation.slides.add_slide(slide_layout)
        #set title
        slide.placeholders[0].text = row[1]
        #set content
        slide.placeholders[1].text = "Duration: " + row[2] + "\nDependencies: " + row[3] + "\nStatus: " + row[4]   + "\nResources: " + row[5] 
    
    presentation.save("project_plan.pptx")
